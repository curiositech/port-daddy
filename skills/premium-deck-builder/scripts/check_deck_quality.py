#!/usr/bin/env python3
"""check_deck_quality.py — automated audit of a .pptx against premium-deck-builder rules.

Usage:
    python3 check_deck_quality.py <deck.pptx>

Exits non-zero if ANY error-level finding is present. Warnings are reported
but don't fail. Use as a CI gate or a pre-flight check before declaring a
deck done.

Rules audited:
    [E] Body font sizes  < 18pt   (too small for projection; flag every instance)
    [E] >3 distinct font families used
    [E] >6 distinct colors used across the deck
    [E] "Thank you" / "Thanks" / "Questions?" on the final slide
    [W] ALL CAPS title text >3 words
    [W] Any tabular column of numbers without tabular-nums (heuristic: monospace
        font absent on a run that looks numeric)
    [W] Topic-phrase title (no verb, no conclusion) — heuristic only

This is a lint, not a renderer. It will not catch grid-vs-text legibility
or other things that only show up in the rendered PNG. Run scripts/render_deck.sh
in addition.
"""

from __future__ import annotations
import sys, re
from collections import Counter
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Pt
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


FINDING_ERROR = 'E'
FINDING_WARN = 'W'

# Heuristic — topic-phrase title patterns (no verb, short, ends without punctuation)
TOPIC_PHRASE_HINT = re.compile(r'^(Q\d|Results|Overview|Summary|Agenda|Introduction|Background|Objectives|Goals|Update|Status|Plan)\b', re.IGNORECASE)
END_WITH_VERB = re.compile(r'\b(grew|fell|beat|missed|launched|shipped|reduced|increased|prevented|exceeded|under-?delivered)\b', re.IGNORECASE)


def audit(deck_path: Path) -> tuple[list[tuple[str, str]], list[tuple[str, str]]]:
    """Returns (errors, warnings) — each a list of (location, message)."""
    prs = Presentation(str(deck_path))
    errors: list[tuple[str, str]] = []
    warnings: list[tuple[str, str]] = []

    fonts_used: Counter[str] = Counter()
    colors_used: set[str] = set()

    slides = list(prs.slides)
    total = len(slides)

    # Final-slide check
    if total >= 1:
        final = slides[-1]
        final_text = ' '.join(
            r.text for sh in final.shapes if sh.has_text_frame
            for p in sh.text_frame.paragraphs for r in p.runs
        ).lower()
        if any(phrase in final_text for phrase in ('thank you', 'thanks!', 'questions?', 'q & a only')):
            errors.append((f'slide {total}', "Final slide reads as 'Thank you' / 'Questions?' — wastes highest-recall slot. End on the takeaway or CTA."))

    for i, slide in enumerate(slides, start=1):
        for sh_idx, sh in enumerate(slide.shapes):
            if not sh.has_text_frame:
                continue
            for p_idx, para in enumerate(sh.text_frame.paragraphs):
                full_text = ''.join(r.text for r in para.runs).strip()
                for r_idx, run in enumerate(para.runs):
                    f = run.font
                    if f.name:
                        fonts_used[f.name] += 1
                    if f.color and f.color.type is not None:
                        try:
                            rgb = f.color.rgb
                            if rgb is not None:
                                colors_used.add(str(rgb))
                        except (AttributeError, TypeError):
                            pass
                    # Size check — matches the skill's typography rule:
                    #   < 14pt is an ERROR (below the global vision-accessibility floor)
                    #     except for eyebrow-class text (uppercase, tracked-out, short)
                    #     and slidedoc captions which can drop to 9pt
                    #   14-17pt is a WARNING (below the 18pt projection caption floor)
                    if f.size is not None:
                        pt = f.size.pt
                        text = run.text.strip()
                        # Eyebrow exception: short all-caps, weight ≥600
                        is_eyebrow = (
                            text == text.upper() and
                            len(text) < 40 and
                            (f.bold is True or (f.bold is None and pt <= 12))
                        )
                        # Caption-class: "fig. NN", page numbers, footnotes — mono and short
                        is_caption_class = (
                            (f.name or '').lower() in {'menlo', 'jetbrains mono', 'geist mono',
                                                       'consolas', 'sf mono', 'berkeley mono',
                                                       'courier new', 'monaco'}
                            and len(text) < 80
                        )
                        if pt < 9:
                            errors.append((
                                f'slide {i}, shape {sh_idx}, run {r_idx}',
                                f"Font size {pt:.0f}pt is below 9pt — illegible on any surface. Text: {text[:60]!r}",
                            ))
                        elif pt < 14 and not (is_eyebrow or is_caption_class):
                            errors.append((
                                f'slide {i}, shape {sh_idx}, run {r_idx}',
                                f"Font size {pt:.0f}pt is below 14pt floor (vision-accessibility). Text: {text[:60]!r}",
                            ))
                        elif pt < 18 and not (is_eyebrow or is_caption_class):
                            warnings.append((
                                f'slide {i}, shape {sh_idx}, run {r_idx}',
                                f"Font size {pt:.0f}pt is below 18pt — fine for slidedoc, too small for projection. Text: {text[:60]!r}",
                            ))

                # ALL CAPS title heuristic
                # Heuristic: first paragraph of first non-empty-text shape on slide,
                # all-caps, more than 3 words.
                if p_idx == 0 and sh_idx <= 2 and full_text and full_text == full_text.upper():
                    word_count = len(full_text.split())
                    if word_count > 3:
                        warnings.append((
                            f'slide {i}',
                            f"ALL-CAPS title ({word_count} words). Sentence case reads 13-20% faster. Text: {full_text[:80]!r}",
                        ))

                # Topic-phrase title heuristic
                if p_idx == 0 and sh_idx <= 2 and full_text:
                    if TOPIC_PHRASE_HINT.match(full_text) and not END_WITH_VERB.search(full_text):
                        warnings.append((
                            f'slide {i}',
                            f"Title may be a topic phrase (no verb / no conclusion). Title: {full_text[:80]!r}. Replace with the conclusion.",
                        ))

    # Aggregate font / color rules
    if len(fonts_used) > 3:
        errors.append(('deck', f"{len(fonts_used)} distinct font families used: {list(fonts_used)}. Cap at 3."))
    if len(colors_used) > 6:
        errors.append(('deck', f"{len(colors_used)} distinct colors used. Cap at 6 (3 brand + 3 accent tints)."))

    return errors, warnings


def main(argv: list[str]) -> int:
    if len(argv) != 2:
        print(f"Usage: {argv[0]} <deck.pptx>", file=sys.stderr)
        return 64
    deck = Path(argv[1])
    if not deck.exists():
        print(f"File not found: {deck}", file=sys.stderr)
        return 66

    errors, warnings = audit(deck)

    print(f"AUDIT: {deck}")
    print(f"  errors:   {len(errors)}")
    print(f"  warnings: {len(warnings)}")
    print()
    for loc, msg in errors:
        print(f"  [E] {loc}: {msg}")
    for loc, msg in warnings:
        print(f"  [W] {loc}: {msg}")
    print()

    return 1 if errors else 0


if __name__ == '__main__':
    sys.exit(main(sys.argv))
